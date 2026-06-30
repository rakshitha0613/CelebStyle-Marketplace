from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

prs = Presentation()

def add_title_slide(title, subtitle=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        try:
            slide.placeholders[1].text = subtitle
        except Exception:
            pass

def add_bullet_slide(title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0

# Title
add_title_slide('CelebStyle — Progress Review', f'Date: {datetime.now().strftime("%Y-%m-%d")})')

# Executive summary
add_bullet_slide('Executive Summary', [
    'Module 3.1 (Admin CMS) completed: full CRUD in frontend and backend',
    'Phase 2 implemented: product pages, cart, checkout (demo), orders, storefronts',
    'Backend: in-memory stores; Frontend: Next.js app router + client components',
])

# Technical inventory
add_bullet_slide('Technical Inventory', [
    'Frontend: Next.js 15, React 19, TypeScript, Tailwind CSS',
    'Backend: Node.js + Express, in-memory stores, REST API (port 4000)',
    'Cart persistence: localStorage (key: celebstyle-cart)',
    'Payment: Razorpay demo simulation (not live)',
])

# Completed features
add_bullet_slide('Completed — Key Features', [
    'Admin CMS: Celebrities, Outfits, Manufacturers CRUD',
    'Product flow: Outfit details, Add-to-cart, Cart page, Checkout demo',
    'Orders: create/list/detail with commission split',
    'Storefronts: builder, listing, commission metrics',
])

# Pending / Caveats
add_bullet_slide('Pending / Caveats', [
    'No persistent DB — orders lost on restart',
    'Payment gateway not integrated (demo only)',
    'Image uploads not implemented (URLs only)',
    'Edge cases and validation need hardening',
])

# Next steps
add_bullet_slide('Recommended Next Steps', [
    'Integrate persistent database (Postgres + Prisma)',
    'Replace Razorpay demo with live integration',
    'Add file upload pipeline for images',
    'Add automated tests and CI checks',
])

# Contact / Notes
add_bullet_slide('Notes', [
    'Dev servers: Frontend (Next dev) on 3003, Backend on 4000',
    'TS check: baseUrl deprecation warning in tsconfig (consider updating)',
    'For demo walkthrough, I can run flows: add-to-cart → checkout → order',
])

output_path = r"C:\Users\Rakshitha\Downloads\CelebStyle_Review_Presentation.pptx"
prs.save(output_path)
print('Saved:', output_path)
